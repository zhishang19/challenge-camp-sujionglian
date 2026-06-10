"""generate_ppt.py — 生成终选答辩 PPT (v2 增强版)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色主题 (专业深蓝+青色) ──
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
MID_BG = RGBColor(0x1B, 0x2A, 0x3A)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2 = RGBColor(0x00, 0x96, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE0, 0xF0, 0xF7)
GRAY = RGBColor(0x90, 0x9A, 0xA3)
GREEN = RGBColor(0x2D, 0xC6, 0x53)
ORANGE = RGBColor(0xFF, 0x9F, 0x1C)
RED = RGBColor(0xE6, 0x3E, 0x3E)
TABLE_HEAD = RGBColor(0x00, 0xB4, 0xD8)
TABLE_ALT = RGBColor(0xE8, 0xF4, 0xFA)
PHASE_BG = {1: RGBColor(0x00, 0x96, 0xC7), 2: RGBColor(0x00, 0xB4, 0x88),
            3: RGBColor(0x00, 0x88, 0xCC), 4: RGBColor(0x7B, 0x2D, 0x8E)}


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_side_bar(slide, color, width=Inches(0.08)):
    """Left accent bar."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), width, prs.slide_height)  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_phase_badge(slide, text, color, left=Inches(0.5), top=Inches(6.6)):
    """Rounded phase badge at bottom."""
    badge = slide.shapes.add_shape(
        5,  # rounded rectangle
        left, top, Inches(2.2), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_slide_number(slide, num, total=16):
    txBox = slide.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(1), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, text, subtitle=None, phase_num=None):
    """Title with optional phase color accent dot."""
    add_side_bar(slide, ACCENT)
    # Phase dot
    if phase_num:
        dot = slide.shapes.add_shape(1, Inches(0.45), Inches(0.38), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = PHASE_BG.get(phase_num, ACCENT)
        dot.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def add_subtitle_line(slide, text, top=Inches(1.05)):
    txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT
    p.font.italic = True


def add_body(slide, text, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), height=Inches(5.5), size=Pt(16), color=RGBColor(0x33, 0x33, 0x44)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(3)


def add_bullets(slide, items, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), size=Pt(15), spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x44)
        p.space_after = spacing
        level = 0
        for ch in item:
            if ch == ' ':
                level += 1
            else:
                break
        p.level = level // 2


def add_table(slide, data, left=Inches(0.8), top=Inches(1.6), col_widths=None, row_height=Inches(0.42)):
    rows = len(data)
    cols = len(data[0])
    w = sum(col_widths) if col_widths else Inches(11.5)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, row_height * rows)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.bold = (r == 0)
                para.font.color.rgb = WHITE if r == 0 else RGBColor(0x33, 0x33, 0x44)
                para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.margin_left = Inches(0.15)
            cell.margin_right = Inches(0.15)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEAD
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT


def add_code_block(slide, code, left=Inches(0.8), top=Inches(1.6), width=Inches(5.5), height=Inches(3.5)):
    """Simulated code block with dark BG."""
    # Background rect
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    rect.line.color.rgb = RGBColor(0x33, 0x33, 0x44)
    rect.line.width = Pt(0.5)
    # Text
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.name = 'Consolas'
        p.font.color.rgb = RGBColor(0xB0, 0xD0, 0xE0)
        p.space_after = Pt(2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: 封面
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
# Decorative top line
line = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
# Big accent rectangle behind title
rect = slide.shapes.add_shape(1, Inches(2.5), Inches(1.8), Inches(8.3), Inches(3.6))
rect.fill.solid(); rect.fill.fore_color.rgb = MID_BG; rect.line.fill.background()
# Title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.3))
tf = txBox.text_frame
p = tf.paragraphs[0]; p.text = "OS Agent 记忆优化"; p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = "及高效应用研究"; p2.font.size = Pt(52); p2.font.bold = True; p2.font.color.rgb = WHITE
# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5))
p = txBox2.text_frame.paragraphs[0]; p.text = "14 天选拔集训 · 成果汇报"; p.font.size = Pt(22); p.font.color.rgb = ACCENT
# Divider
div = slide.shapes.add_shape(1, Inches(0.8), Inches(4.4), Inches(2.5), Inches(0.04))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
# Info
txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.5), Inches(1))
tf = txBox3.text_frame
p = tf.paragraphs[0]; p.text = "苏炯濂"; p.font.size = Pt(20); p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = "2026-06-09"; p2.font.size = Pt(14); p2.font.color.rgb = GRAY
# Bottom bar
bot = slide.shapes.add_shape(1, Inches(0), Inches(7.2), prs.slide_width, Inches(0.3))
bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()
add_slide_number(slide, 1, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: 目录
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "目  录", "CONTENTS")
add_subtitle_line(slide, "——————————————————————————————————————————")
toc = [
    "01    赛题理解与集训定位", "",
    "02    项目技术架构与数据流", "",
    "03    Phase 1 · 数据清洗与 Git 基础", "",
    "04    Phase 2 · 多源数据整合与质量评测", "",
    "05    Phase 3 · AI 应用开发", "",
    "06    Phase 4 · 麒麟环境与数据库", "",
    "07    成果数据统计", "",
    "08    技能收获与展望",
]
add_bullets(slide, toc, left=Inches(1.5), top=Inches(1.5), size=Pt(18), spacing=Pt(5))
add_slide_number(slide, 2, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: 赛题理解
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "赛题理解与集训定位", "OS Agent 记忆优化（偏好记忆 + 知识记忆）", 1)
add_subtitle_line(slide, "——————————————————————————————————————————")
# Left column: 赛题要求
p1 = [
    "▎ 赛题核心能力",
    "  • 多源数据整合、清洗、格式标准化与质量校验",
    "  • 偏好记忆的动态捕捉与版本化管理",
    "  • 知识记忆的结构化整合与关联检索（端侧 ≤500ms）",
    "  • 敏感信息识别与过滤（PII 脱敏）",
    "  • 短期 / 中期 / 长期记忆数据流转",
]
add_bullets(slide, p1, left=Inches(0.8), top=Inches(1.5), width=Inches(6.0), size=Pt(14))
# Right column: 集训信息
p2 = [
    "▎ 集训信息",
    "  • 性质：选拔性集训 | 14 天 | 约 25 人",
    "  • 选拔：Top 10 入队 + 5 替补",
    "  • 形式：资料自学 + 教师答疑",
    "  • 项目：个人独立仓库",
    "  • 评测：四阶段逐级淘汰",
]
add_bullets(slide, p2, left=Inches(7.2), top=Inches(1.5), width=Inches(5.5), size=Pt(14))
# Bottom: 阶段晋级表
add_table(slide, [
    ["阶段", "建议天数", "晋级规则", "我的状态"],
    ["Phase 1 基础", "1-4 天", "全员通过", "已完成 (Git 3 次提交)"],
    ["Phase 2 巩固", "3-5 天", "≥55 分，累计 ≥120", "已完成 (报告+数据)"],
    ["Phase 3 提升", "3-5 天", "≥60 分，累计 ≥175", "已完成 (5/5 测例)"],
    ["Phase 4 拓展", "5 天", "Top 10 入队", "已完成 (PPT+SQLite)"],
], top=Inches(4.2), col_widths=[Inches(2.2), Inches(2.2), Inches(3.5), Inches(3.5)], row_height=Inches(0.42))
add_slide_number(slide, 3, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: 技术架构
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "项目技术架构与数据流", "4 层架构 · Python + Java + SQLite + Git", 2)
add_subtitle_line(slide, "——————————————————————————————————————————")
# Left: 数据流图（ASCII art style in code block）
arch_text = (
    "┌── 数据层 ─────────────────────────┐\n"
    "│ raw/d2,d3  → merge_day2 / clean_basic   │\n"
    "│ user_behavior, tool_result, chat_sessions │\n"
    "│ 输出: merged.jsonl (9行), clean.csv        │\n"
    "└────────────────────────────────────┘\n"
    "               ↓\n"
    "┌── 整合层 (Phase 2) ───────────────┐\n"
    "│ consolidate.py → 7 模块 JSON         │\n"
    "│ chat_turns / preferences / knowledge   │\n"
    "│ tool_exec / memory_events / snapshots   │\n"
    "│ + quality_evaluation (10/10 hit)        │\n"
    "└────────────────────────────────────┘\n"
    "               ↓\n"
    "┌── 应用层 (Phase 3) ───────────────┐\n"
    "│ LangChain4j /api/chat + LangGraph      │\n"
    "│ Mini Memory Agent CLI + TEST.md        │\n"
    "│ 5/5 test cases passed                  │\n"
    "└────────────────────────────────────┘\n"
    "               ↓\n"
    "┌── 存储层 (Phase 4) ───────────────┐\n"
    "│ SQLite memory.db (5 tables, 36KB)      │\n"
    "│ 5 CSV exports + Benchmark (0.33ms avg) │\n"
    "└────────────────────────────────────┘"
)
add_code_block(slide, arch_text, left=Inches(0.8), top=Inches(1.5), width=Inches(6.2), height=Inches(5.5))
# Right: 技术栈
p = [
    "▎ 技术栈",
    "  Python 3.12+    数据清洗、整合、Agent",
    "  Java 17          Spring Boot 3.3.5",
    "  LangChain4j      LLM 对话接口",
    "  LangGraph        4 节点工作流",
    "  SQLite           端侧轻量数据库",
    "  Git / GitHub     4 分支版本管理",
    "",
    "▎ 核心指标",
    "  检索延迟  0.33ms avg (目标 ≤500ms)",
    "  PII 脱敏  100%",
    "  质量评测  10/10 hit",
    "  测试通过  5/5 tests",
    "  数据库    memory.db 36KB",
    "  CSV 导出  5 文件",
    "  录屏      demo.mp4 21MB",
]
add_bullets(slide, p, left=Inches(7.3), top=Inches(1.5), width=Inches(5.5), size=Pt(12.5), spacing=Pt(3))
add_slide_number(slide, 4, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Phase 1
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Phase 1 · 数据清洗与 Git 基础", "D2 多源合并 / D3 CSV 清洗 / 分支合并含冲突处理", 1)
add_subtitle_line(slide, "——————————————————————————————————————————")
add_phase_badge(slide, "PHASE 1 — 基础阶段", PHASE_BG[1])
# Left: D2 + D3
p1 = [
    "▎ D2 · merge_day2.py",
    "  读取 raw/d2/user_behavior.json (8条) + tool_result.json (3条)",
    "  处理流程:",
    "    ① 字段校验 (uid/action/content 必填, status 枚举校验)",
    "    ② 文本清洗 (HTML 标签、emoji、填充词、错别字修正)",
    "    ③ 时间标准化 (8 种格式 → YYYY-MM-DD HH:MM:SS)",
    "    ④ 大颗粒度去重 (uid+action+content 三元组)",
    "  输出: merged.jsonl (9 行，去重 1 条重复)",
    "",
    "▎ D3 · clean_basic.py",
    "  读取 chat_sessions_dirty.csv",
    "  清洗: HTML 清理 / 填充词移除 / 手机号邮箱脱敏",
    "  输出: chat_sessions_clean.csv + review_log.txt",
]
add_bullets(slide, p1, left=Inches(0.8), top=Inches(1.5), width=Inches(6.5), size=Pt(12.5), spacing=Pt(3))
# Right: Git + 数据示例
p2 = [
    "▎ Git · do_git_merge.py",
    "  演示含冲突的 Git 分支合并流程:",
    "  • git init → 创建 main + feature 分支",
    "  • 模拟 3 类冲突: clear cut / conflict / auto merge",
    "  • 解决策略: keep_latest",
    "",
    "▎ 验收清单",
    "  AI 工具 >=2        Trae AI + 通义千问",
    "  README.md         环境清单 (OS/Python/Git/IDE)",
    "  Git 提交 >=2       3 次提交 (phase1-basics)",
    "  .gitignore        敏感文件排除",
    "  clean.log         字段校验 + 清洗操作日志",
]
add_bullets(slide, p2, left=Inches(7.5), top=Inches(1.5), width=Inches(5.3), size=Pt(12.5), spacing=Pt(3))
add_slide_number(slide, 5, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Phase 2 (Part 1)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Phase 2 · 多源数据整合", "consolidate.py — 7 大模块 / 6 类 JSON / 质量评测", 2)
add_subtitle_line(slide, "——————————————————————————————————————————")
add_phase_badge(slide, "PHASE 2 — 巩固阶段", PHASE_BG[2])
add_table(slide, [
    ["模块", "JSON 文件", "记录数", "关键处理"],
    ["chat_turns", "chat_turns.json", "7 条", "uid/role/time 标准化, PII 脱敏, 去重"],
    ["preferences", "preferences.json", "4 条", "2 冲突: output_style / emoji_style, keep_latest"],
    ["knowledge_items", "knowledge_items.json", "7 条", "标题/正文/标签/来源/类型, 过滤无效案例"],
    ["tool_executions", "tool_executions.json", "0 条", "D4 tool_log 过滤 (非可执行记录)"],
    ["memory_events", "memory_events.json", "14 条", "7 chat + 7 knowledge items"],
    ["memory_snapshots", "memory_snapshots_resolved.json", "6 条", "3 冲突已解决 (output_style/emoji/answer)"],
    ["quality_evaluation", "quality_evaluation.json", "10 案例", "10/10 hit, PII leak 0/10"],
], top=Inches(1.5), col_widths=[Inches(2.0), Inches(3.5), Inches(1.5), Inches(4.5)])
# Bottom: 质量 + 报告
p = [
    "▎ 核心机制    PII 脱敏: 邮箱 → [REDACTED], 手机号 → [REDACTED]　|　冲突合并: keep_latest / keep_previous / needs_review",
    "▎ 质量报告    report.md (清洗规则/异常计数/样例对比/停顿词统计) + quality_eval_report.md (10 案例逐条分析)"
]
add_bullets(slide, p, left=Inches(0.8), top=Inches(5.3), width=Inches(11.5), size=Pt(12), spacing=Pt(3))
add_slide_number(slide, 6, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Phase 2 (Part 2 — 冲突与质量)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Phase 2 · 冲突解决与质量评测", "3 类偏好冲突 / 10 个质量案例 / 100% 命中率", 2)
add_subtitle_line(slide, "——————————————————————————————————————————")
add_phase_badge(slide, "PHASE 2 — 巩固阶段", PHASE_BG[2])
# 冲突表
add_table(slide, [
    ["用户", "冲突键", "版本 v1", "版本 v2", "胜出者", "策略"],
    ["u001", "output_style", "详细版+表格 (v1)", "简洁版 (v2)", "简洁版", "keep_latest"],
    ["u002", "emoji_style", "禁用 emoji (v1)", "允许 emoji (v2)", "允许 emoji", "keep_latest"],
    ["u007", "answer_style", "先结论再步骤 (v1)", "详细解释 (v2)", "详细解释", "keep_latest"],
], top=Inches(1.5), col_widths=[Inches(1.2), Inches(2.0), Inches(2.5), Inches(2.5), Inches(1.8), Inches(1.5)])
# 质量评测
add_table(slide, [
    ["编号", "查询", "状态", "PII 泄露", "分析"],
    ["Q001", "帮我导出月报，按我的偏好来", "hit", "否", "偏好感知查询 → 知识检索"],
    ["Q003", "麒麟系统怎么更新驱动？", "hit", "否", "直接命中 K004 (麒麟装驱动)"],
    ["Q005", "把今天会议整理一下", "hit", "否", "匹配 K003 + 同义词扩展"],
    ["Q008", "离线安装 deb 包失败怎么办？", "hit", "否", "deb→dpkg 同义词扩展命中"],
    ["Q010", "web_search 超时，驱动知识不能用？", "hit", "否", "工具超时 + 知识回退"],
    ["汇总", f"10 案例全部命中", "100%", "0/10", "hit_rate=100%, pii_leak=0"],
], top=Inches(3.4), col_widths=[Inches(1.0), Inches(3.5), Inches(1.0), Inches(1.2), Inches(4.8)], row_height=Inches(0.4))
add_slide_number(slide, 7, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Phase 3 (Part 1)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Phase 3 · LangChain4j + LangGraph", "Java Spring Boot API / Python 4 节点工作流", 3)
add_subtitle_line(slide, "——————————————————————————————————————————")
add_phase_badge(slide, "PHASE 3 — 提升阶段", PHASE_BG[3])
# Left: LangChain4j
p1 = [
    "▎ LangChain4j (Java Spring Boot 3.3.5)",
    "  • GET /api/chat?q=xxx&uid=u001  — LLM 对话接口",
    "  • PreferenceLoader 读取 preferences.json → SystemMessage",
    "  • 支持 Ollama / OpenAI 切换 (CAMP_LLM_PROVIDER 环境变量)",
    "  • API Key 不入库 — .env.example 替代，.env 在 .gitignore",
    "  • mvn compile 通过，5 个 Java 源文件",
    "",
    "▎ LangGraph 工作流 (Pure Python)",
    "  • 4 节点:",
    "    ① filter_sensitive  — PII 模式检测 (email/phone/id_card/ip)",
    "    ② classify          — 意图分类 (knowledge/casual/disable)",
    "    ③ retrieve / chat    — Bigram 检索 / 普通回复",
    "    ④ output             — 格式化 + 偏好注入",
    "  • 知识源: Phase 2 knowledge_items.json (7 条)",
]
add_bullets(slide, p1, left=Inches(0.8), top=Inches(1.5), width=Inches(6.3), size=Pt(12.5), spacing=Pt(2.5))
# Right: code + diagram
code = (
    "# LangGraph 4-node pipeline\n"
    "state = {'query': user_input}\n"
    "state = filter_sensitive(state)\n"
    "state = classify(state)\n"
    "if state['intent'] == 'know':\n"
    "    state = retrieve(state)\n"
    "else:\n"
    "    state = chat(state)\n"
    "state = output(state)"
)
add_code_block(slide, code, left=Inches(7.3), top=Inches(1.5), width=Inches(5.5), height=Inches(2.6))
p2 = [
    "▎ 技术要点",
    "  • Bigram 检索: O(n) 纯 Python, 无需 embedding",
    "  • 同义词扩展: 月报→会议纪要, deb→dpkg",
    "  • 意图分类: 关键词匹配 (不调 LLM)",
    "  • PII 匹配: 邮箱/手机号/身份证/IP 4 种模式",
]
add_bullets(slide, p2, left=Inches(7.3), top=Inches(4.3), width=Inches(5.5), size=Pt(12.5), spacing=Pt(2.5))
add_slide_number(slide, 8, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Phase 3 (Part 2 — Mini Agent)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Phase 3 · Mini Memory Agent", "CLI 交互式 Agent · 5/5 测例通过 · 录屏 demo.mp4", 3)
add_subtitle_line(slide, "——————————————————————————————————————————")
add_phase_badge(slide, "PHASE 3 — 提升阶段", PHASE_BG[3])
add_table(slide, [
    ["测试", "用户", "输入", "预期行为", "验证点", "结果"],
    ["Test 1", "u001", "麒麟系统如何更新驱动", "检索 K004+K005", "bigram 匹配 + 偏好注入", "PASS"],
    ["Test 2", "u001", "写月报摘要", "检索 K003", "月报→会议纪要 同义词扩展", "PASS"],
    ["Test 3", "u002", "你好", "非知识查询 → 记忆快照", "emoji 偏好检测", "PASS"],
    ["Test 4", "u003", "离线安装 .deb 包", "检索 K002+dpkg", "deb→dpkg 同义词扩展", "PASS"],
    ["Test 5", "u004", "liubei@shu.com 别记下", "PII 脱敏 + forget", "email redact + forget 标记", "PASS"],
], top=Inches(1.5), col_widths=[Inches(1.0), Inches(1.0), Inches(2.8), Inches(2.5), Inches(2.5), Inches(1.0)], row_height=Inches(0.5))
p = [
    "▎ Agent 架构   输入 → PII 检测 → 意图分类 → 知识检索 (bigram+同义词) → 偏好注入 → 记忆快照更新 → 输出",
    "▎ 数据来源     knowledge_items.json (7条) + preferences.json (4条) + memory_snapshots.json (6条)",
    "▎ 运行方式     python mini_memory_agent.py --test  (5 测例自动验证)   |   python mini_memory_agent.py  (交互模式)",
    "▎ 交付物       TEST.md (5 条测例) + demo.mp4 (21MB 录屏) + min_memory_agent.py (500+ 行)",
]
add_bullets(slide, p, left=Inches(0.8), top=Inches(5.3), width=Inches(11.5), size=Pt(12), spacing=Pt(3))
add_slide_number(slide, 9, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Phase 4
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Phase 4 · 麒麟环境与数据库", "SQLite 5 表 · CSV 导出 · 检索基准 ≤500ms", 4)
add_subtitle_line(slide, "——————————————————————————————————————————")
add_phase_badge(slide, "PHASE 4 — 拓展阶段", PHASE_BG[4])
# SQLite 表结构
add_table(slide, [
    ["表名", "字段", "记录", "说明"],
    ["knowledge_items", "id, title, body, tags, type, source", "7", "知识条目 (含 K001-K007)"],
    ["preferences", "uid, pref_key, pref_value, type, version, ttl", "4", "用户偏好 (2 冲突已解决)"],
    ["chat_turns", "session_id, user_id, role, message, created_at", "7", "对话历史 (PII 已脱敏)"],
    ["tool_executions", "tool_id, source, content, status, duration_ms", "0", "工具调用记录 (全部 tool_log)"],
    ["memory_snapshots", "uid, key, memory_value, version, scope", "6", "记忆快照 (3 冲突 resolved)"],
], top=Inches(1.5), col_widths=[Inches(2.2), Inches(4.5), Inches(1.2), Inches(3.5)])
# Benchmark 数据
p = [
    "▎ 麒麟环境   kylin_setup.md — 安装流程 / 基础命令 / 环境配置 / 排查指南",
    "▎ CSV 导出   exports/ 目录 — 5 个 CSV (utf-8-sig), Excel 直接打开无乱码",
    "▎ 检索基准   benchmark_retrieval.py — 750 样本 (15 query x 50 iterations)"
]
add_bullets(slide, p, left=Inches(0.8), top=Inches(4.2), width=Inches(7.5), size=Pt(12), spacing=Pt(3))
add_table(slide, [
    ["指标", "Avg", "Median (P50)", "P95", "P99", "Min", "Max", "Pass Rate"],
    ["延迟", "0.33ms", "0.10ms", "1.02ms", "1.06ms", "0.05ms", "1.23ms", "100% (750/750)"],
], top=Inches(5.5), col_widths=[Inches(1.3), Inches(1.3), Inches(1.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(2.0)], row_height=Inches(0.42))
add_slide_number(slide, 10, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Phase 4 (Part 2 — PPT & Q&A)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Phase 4 · 终选答辩准备", "PPT 制作 + 理论 Q&A 33 题", 4)
add_subtitle_line(slide, "——————————————————————————————————————————")
add_phase_badge(slide, "PHASE 4 — 拓展阶段", PHASE_BG[4])
p1 = [
    "▎ 答辩 PPT (15 页)",
    "  封面 → 目录 → 赛题理解 → 技术架构 → Phase 1-4 (各 1-2 页) → 成果统计 → 展望 → 致谢",
    "  设计: 深蓝 + 青色专业配色 / 侧边栏导航 / 数据驱动 / 代码片段",
    "",
    "▎ 理论准备 (答辩理论Q_A.md — 33 题)",
    "  数据清洗 (Q1-Q4)     PII、去重、时间标准化、停顿词",
    "  Git 版本管理 (Q5-Q8)  四区模型、冲突解决、.gitignore",
    "  Python 基础 (Q9-Q12)  列表/字典、__main__、JSON、正则",
    "  AI 应用 (Q13-Q19)     LangChain4j、LangGraph、bigram、同义词扩展",
    "  数据库 (Q20-Q22)      SQLite 选型、5 张表、CSV 导出",
    "  麒麟 OS (Q23-Q25)     麒麟 vs Ubuntu、安装命令",
    "  安全隐私 (Q26-Q27)    PII 脱敏、forget 机制、API Key",
    "  工程素养 (Q28-Q30)    requirements.txt、README、收获",
    "  Mini Agent (Q31-Q33)  架构、5 测例、≤500ms 原理",
]
add_bullets(slide, p1, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), size=Pt(13), spacing=Pt(2.5))
add_slide_number(slide, 11, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: 成果数据
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "成果数据总览", "4 个阶段 · 输出 6 类 JSON · 5 个 CSV · 100% 测试通过")
add_subtitle_line(slide, "——————————————————————————————————————————")
add_table(slide, [
    ["维度", "指标", "数值", "状态"],
    ["Phase 1 数据", "merged.jsonl (去重后)", "9 行 (>=8 要求)", "通过"],
    ["Phase 1 数据", "clean_basic CSV 清洗", "chat_sessions_clean.csv", "通过"],
    ["Phase 1 Git", "Git 提交次数", "3 次 (phase1-basics)", "通过"],
    ["Phase 2 整合", "JSON 输出", "6 类文件 (24 条记录)", "通过"],
    ["Phase 2 质量", "Quality Eval", "10/10 hit, 0 PII leak", "通过"],
    ["Phase 2 冲突", "偏好冲突解决", "3 冲突 (keep_latest)", "通过"],
    ["Phase 3 Java", "Spring Boot 编译", "mvn compile 通过", "通过"],
    ["Phase 3 Agent", "Mini Agent 测例", "5/5 PASS", "通过"],
    ["Phase 3 录屏", "demo.mp4", "21 MB", "通过"],
    ["Phase 4 SQLite", "memory.db", "36 KB, 5 tables", "通过"],
    ["Phase 4 导出", "CSV", "5 files (utf-8-sig)", "通过"],
    ["Phase 4 检索", "Benchmark (750 样本)", "avg 0.33ms, 100% ≤500ms", "通过"],
    ["Phase 4 PPT", "终选答辩_v2.pptx", "15 slides", "通过"],
    ["Phase 4 理论", "答辩理论Q_A.md", "33 题覆盖", "通过"],
], top=Inches(1.3), col_widths=[Inches(1.8), Inches(2.8), Inches(3.5), Inches(1.5)], row_height=Inches(0.35))
add_slide_number(slide, 12, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: Git 与工程规范
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "Git 版本管理 & 工程规范", "4 分支 / 5 次提交 / 仓库地址: github.com/zhishang19/challenge-camp-sujionglian")
add_subtitle_line(slide, "——————————————————————————————————————————")
# Git branch structure
add_table(slide, [
    ["分支", "提交数", "关键文件", "验收点"],
    ["phase1-basics", "3 次", "merge_day2 / clean_basic / do_git_merge + README", "AI工具>=2, Git>=2次, .gitignore"],
    ["phase2-consolidate", "2 次", "consolidate.py + 6 JSON + report.md + requirements.txt", "4类JSON+报告, Git笔记, pipeline"],
    ["phase3-advance", "2 次", "LangChain4j + LangGraph + MiniAgent + TEST.md + demo.mp4", "/api/chat, 工作流, 5测例, 录屏"],
    ["phase4-extend", "1 次", "setup_db.py + benchmark + kylin_setup + PPT + Q&A", "SQLite, CSV导出, PPT, 检索基准"],
], top=Inches(1.5), col_widths=[Inches(2.2), Inches(1.2), Inches(5.5), Inches(2.5)])
p = [
    "▎ 工程规范    .gitignore (.env/__pycache__/target/) + requirements.txt + D05-Git协作.md (Git 使用笔记)",
    "▎ 安全措施    API Key 不入库 (.env.example) · 敏感信息 .gitignore 排除 · demo.mp4 录屏证明可运行 · TEST.md 自动化验证"
]
add_bullets(slide, p, left=Inches(0.8), top=Inches(4.2), width=Inches(11.5), size=Pt(13), spacing=Pt(3))
add_slide_number(slide, 13, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14: 技能收获
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title(slide, "技能收获与成长", "5 大能力域 · 从零到交付完整 Agent 记忆管线")
add_subtitle_line(slide, "——————————————————————————————————————————")
skills = [
    ("Python 数据工程", "JSON/CSV/JSONL 处理 → 正则表达式 → 文本清洗 → dedup → 时间标准化 → 日志记录"),
    ("Git 版本管理", "4 分支独立开发 → merge 冲突解决 → .gitignore 敏感文件排除 → commit 规范"),
    ("AI 应用开发", "LangChain4j Spring Boot → /api/chat 接口 → SystemMessage 偏好注入 → LangGraph 4 节点工作流 → bigram 检索 → 同义词扩展 → PII 脱敏"),
    ("数据库", "SQLite 建表/导入/查询/导出 → CSV 编码 (utf-8-sig) → 5 表 schema 设计"),
    ("国产 OS 意识", "麒麟操作系统环境认知 → 安装流程 → apt 包管理 → 端侧部署概念"),
]
y = Inches(1.5)
for title, desc in skills:
    # Skill title box
    box = slide.shapes.add_shape(1, Inches(0.8), y, Inches(2.0), Inches(0.55))
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT; box.line.fill.background()
    tf = box.text_frame; p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    # Description
    txBox = slide.shapes.add_textbox(Inches(3.0), y, Inches(9.5), Inches(0.55))
    p = txBox.text_frame.paragraphs[0]; p.text = desc
    p.font.size = Pt(12); p.font.color.rgb = RGBColor(0x33, 0x33, 0x44)
    y += Inches(0.65)

add_slide_number(slide, 14, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15: 展望 & Thank you
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
# Left content
p1 = [
    "▎ 短期改进",
    "  • 接入麒麟 embedding SDK 替代 bigram",
    "  • 实现端侧向量数据库 (ChromaDB/LanceDB)",
    "  • 多轮对话上下文管理",
    "  • 冲突 Explainability (A/B 版本对比)",
    "",
    "▎ 长期方向",
    "  • 完整 Agent 记忆模块架构",
    "  • 麒麟 OS 端侧部署与性能测试",
    "  • 安全审计 (PII 全链路追踪)",
]
add_bullets(slide, p1, left=Inches(0.8), top=Inches(1.5), width=Inches(6.0), size=Pt(15), spacing=Pt(4))
# Right: Thank you
add_body(slide, "谢 谢", left=Inches(7.5), top=Inches(2.5), width=Inches(5.5), height=Inches(1.5), size=Pt(64), color=WHITE)
add_body(slide, "THANK YOU", left=Inches(7.5), top=Inches(3.8), width=Inches(5.5), height=Inches(1.0), size=Pt(32), color=ACCENT)
info = [
    "GitHub: github.com/zhishang19/challenge-camp-sujionglian",
    "",
    "Phase 1 · 基础  →  3 commits, merged.jsonl 9 rows",
    "Phase 2 · 巩固  →  6 JSON, 10/10 quality hit",
    "Phase 3 · 提升  →  5/5 tests, demo.mp4 21MB",
    "Phase 4 · 拓展  →  SQLite 36KB, PPT 15 slides",
]
add_bullets(slide, info, left=Inches(7.5), top=Inches(5.0), width=Inches(5.5), size=Pt(12), spacing=Pt(3))
# Bottom bar
bot = slide.shapes.add_shape(1, Inches(0), Inches(7.2), prs.slide_width, Inches(0.3))
bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()
add_slide_number(slide, 15, 15)

# ── 保存 ──
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "终选答辩_v2.pptx")
prs.save(OUTPUT)
print(f"PPT saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
